"""
Generate Kaigo Navigator client demo presentation — clean light theme.
Run: python scripts/generate_ppt.py
Output: docs/kaigo_navigator_demo.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colour palette (clean enterprise / healthcare) ─────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)   # headings, hero bg
TEAL       = RGBColor(0x00, 0xB3, 0x88)   # primary accent
SOFT_GRAY  = RGBColor(0xF4, 0xF6, 0xF9)   # card backgrounds
MID_GRAY   = RGBColor(0xE2, 0xE8, 0xF0)   # dividers, alt rows
SLATE      = RGBColor(0x4A, 0x55, 0x68)   # body text
SUBTEXT    = RGBColor(0x71, 0x80, 0x96)   # secondary text
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)   # in-progress / warning
MUTED_RED  = RGBColor(0xE5, 0x3E, 0x3E)   # danger / "without"
TEAL_DARK  = RGBColor(0x00, 0x7A, 0x5E)   # deep teal for "with" panel
TEAL_LIGHT = RGBColor(0xE6, 0xF7, 0xF3)   # very light teal for success bg

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color: RGBColor = WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color,
             line_color=None, line_width_pt=None):
    from pptx.util import Pt as PtU
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width_pt:
            shape.line.width = PtU(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=14, bold=False, color=SLATE,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def section_header(slide, title, subtitle=None):
    """Shared top header: teal bar + navy title + optional subtitle."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), TEAL)
    add_text(slide, title,
             Inches(0.6), Inches(0.22), Inches(10), Inches(0.65),
             font_size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.6), Inches(0.88), Inches(12.1), Inches(0.42),
                 font_size=14, color=SUBTEXT, italic=True)
    # thin divider under header
    add_rect(slide, Inches(0.6), Inches(1.12), Inches(12.1), Inches(0.02), MID_GRAY)


def card(slide, left, top, width, height, accent_left=True):
    """White card with soft shadow simulation and optional teal left bar."""
    add_rect(slide, left, top, width, height, SOFT_GRAY,
             line_color=MID_GRAY, line_width_pt=0.5)
    if accent_left:
        add_rect(slide, left, top, Inches(0.07), height, TEAL)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)

    # Hero navy block (top 55%)
    add_rect(slide, 0, 0, SLIDE_W, Inches(4.1), NAVY)

    # Teal accent bar at bottom of hero
    add_rect(slide, 0, Inches(4.1), SLIDE_W, Inches(0.08), TEAL)

    # Japanese title
    add_text(slide, "介護ナビゲーター",
             Inches(1), Inches(0.7), Inches(11.3), Inches(1.1),
             font_size=54, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # English title
    add_text(slide, "Kaigo Navigator",
             Inches(1), Inches(1.75), Inches(11.3), Inches(0.8),
             font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Tagline
    add_text(slide, "Autonomous Multi-Agent AI for Japan's Eldercare Crisis",
             Inches(1), Inches(2.55), Inches(11.3), Inches(0.5),
             font_size=17, color=MID_GRAY, align=PP_ALIGN.CENTER, italic=True)

    # Thin teal divider
    add_rect(slide, Inches(5.0), Inches(3.2), Inches(3.3), Inches(0.04), TEAL)

    # Sub-tagline
    add_text(slide, "One request  ·  Five AI agents  ·  Full care plan in 15 seconds",
             Inches(1), Inches(3.35), Inches(11.3), Inches(0.45),
             font_size=15, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Three stat pills (white bg section)
    pill_data = [
        ("15 sec", "End-to-end pipeline"),
        ("5 agents", "Working in parallel"),
        ("3 forms", "Pre-filled by AI"),
    ]
    pw = Inches(3.3)
    ph = Inches(1.5)
    pg = Inches(0.41)
    px = Inches(0.92)
    py = Inches(4.55)

    for num, label in pill_data:
        add_rect(slide, px, py, pw, ph, WHITE,
                 line_color=MID_GRAY, line_width_pt=0.5)
        add_rect(slide, px, py, pw, Inches(0.05), TEAL)
        add_text(slide, num, px, py + Inches(0.15), pw, Inches(0.65),
                 font_size=30, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        add_text(slide, label, px, py + Inches(0.82), pw, Inches(0.45),
                 font_size=12, color=SLATE, align=PP_ALIGN.CENTER)
        px += pw + pg

    # Footer
    add_rect(slide, 0, SLIDE_H - Inches(0.45), SLIDE_W, Inches(0.45), SOFT_GRAY)
    add_text(slide, "Built by Ganesh Suni Jathar",
             Inches(0.5), SLIDE_H - Inches(0.42), Inches(7), Inches(0.38),
             font_size=11, color=SUBTEXT)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
def slide_problem(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    section_header(slide, "The Problem",
                   "Japan is aging faster than any country in history — the care system hasn't kept up.")

    stats = [
        ("30%",        "Population\nover 65", "Highest globally"),
        ("11 Million", "Care worker\nshortage", "Projected by 2040"),
        ("4–6 Weeks",  "To complete\n介護保険 application", "Per patient, manually"),
        ("3–7 Forms",  "Per patient",  "All Japanese bureaucracy"),
    ]

    bw = Inches(2.8)
    bh = Inches(2.5)
    bg = Inches(0.37)
    bx = Inches(0.57)
    by = Inches(1.35)

    for num, label, sub in stats:
        add_rect(slide, bx, by, bw, bh, WHITE,
                 line_color=MID_GRAY, line_width_pt=0.75)
        add_rect(slide, bx, by, bw, Inches(0.07), TEAL)
        add_text(slide, num, bx, by + Inches(0.18), bw, Inches(0.85),
                 font_size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, label, bx, by + Inches(1.0), bw, Inches(0.75),
                 font_size=13, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
        add_text(slide, sub, bx, by + Inches(1.75), bw, Inches(0.5),
                 font_size=11, color=SUBTEXT, align=PP_ALIGN.CENTER)
        bx += bw + bg

    # Body text
    add_text(slide,
             "Families spend weeks manually matching facilities, filling kanji-dense government "
             "forms, and coordinating visit schedules — often with no professional guidance. "
             "Coordinators are overwhelmed. Facilities go unmatched. Patients wait.",
             Inches(0.6), Inches(4.1), Inches(12.1), Inches(0.85),
             font_size=13, color=SLATE)

    # Solution callout
    add_rect(slide, Inches(0.6), Inches(5.1), Inches(12.1), Inches(0.88), TEAL_LIGHT,
             line_color=TEAL, line_width_pt=1.0)
    add_rect(slide, Inches(0.6), Inches(5.1), Inches(0.07), Inches(0.88), TEAL)
    add_text(slide,
             "💡  Kaigo Navigator replaces weeks of manual work with a 15-second AI pipeline.",
             Inches(0.82), Inches(5.2), Inches(11.7), Inches(0.65),
             font_size=15, bold=True, color=TEAL_DARK)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — Solution Overview
# ══════════════════════════════════════════════════════════════════════════════
def slide_solution(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    section_header(slide, "The Solution",
                   "5 AI Agents · One Pipeline · Human Always in Control")

    agents = [
        ("01", "Orchestrator",      "Routes requests,\nmanages LangGraph state",     "Llama 3.3 70B"),
        ("02", "Service Discovery", "Maps needs → 介護保険 codes\nRAG over facility DB","Llama 3.3 70B + RAG"),
        ("03", "Paperwork Agent",   "Pre-fills all required\ngovernment forms",        "Llama 3.1 8B"),
        ("04", "Scheduling Agent",  "4-week visit calendar\n+ LINE reminders",         "Rule-based"),
        ("05", "Monitoring Agent",  "Scores care plan 0–100\nSurfaces risk alerts",    "Llama 3.1 8B"),
    ]

    aw = Inches(2.32)
    ah = Inches(3.1)
    ag = Inches(0.21)
    ax = Inches(0.43)
    ay = Inches(1.35)

    for num, name, desc, model in agents:
        add_rect(slide, ax, ay, aw, ah, WHITE,
                 line_color=MID_GRAY, line_width_pt=0.75)
        add_rect(slide, ax, ay, aw, Inches(0.07), TEAL)

        # Number badge
        add_rect(slide, ax + Inches(0.15), ay + Inches(0.15),
                 Inches(0.52), Inches(0.3), TEAL)
        add_text(slide, num,
                 ax + Inches(0.15), ay + Inches(0.13),
                 Inches(0.52), Inches(0.3),
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_text(slide, name,
                 ax + Inches(0.12), ay + Inches(0.55),
                 aw - Inches(0.24), Inches(0.52),
                 font_size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, desc,
                 ax + Inches(0.12), ay + Inches(1.1),
                 aw - Inches(0.24), Inches(1.15),
                 font_size=11, color=SLATE, align=PP_ALIGN.CENTER)

        # Model pill
        add_rect(slide, ax + Inches(0.18), ay + ah - Inches(0.52),
                 aw - Inches(0.36), Inches(0.38), TEAL_LIGHT,
                 line_color=TEAL, line_width_pt=0.5)
        add_text(slide, model,
                 ax + Inches(0.18), ay + ah - Inches(0.52),
                 aw - Inches(0.36), Inches(0.38),
                 font_size=10, color=TEAL_DARK, align=PP_ALIGN.CENTER)

        ax += aw + ag

    # HITL banner
    add_rect(slide, Inches(0.43), Inches(4.65), Inches(12.47), Inches(0.72),
             RGBColor(0xFF, 0xF8, 0xE6), line_color=AMBER, line_width_pt=1.0)
    add_rect(slide, Inches(0.43), Inches(4.65), Inches(0.07), Inches(0.72), AMBER)
    add_text(slide,
             "⏸  Human-in-the-loop gate fires before any form submission — "
             "coordinator reviews and approves with one click.",
             Inches(0.65), Inches(4.73), Inches(12.1), Inches(0.55),
             font_size=13, bold=True, color=RGBColor(0x92, 0x40, 0x00))

    # Tech badges
    badges = ["LangGraph 0.6", "FastAPI", "Next.js 16", "Pinecone", "Groq (Llama)", "Langfuse"]
    bx = Inches(0.43)
    for b in badges:
        bw = Inches(1.88)
        add_rect(slide, bx, Inches(5.6), bw, Inches(0.36), SOFT_GRAY,
                 line_color=MID_GRAY, line_width_pt=0.5)
        add_text(slide, b, bx + Inches(0.05), Inches(5.62),
                 bw - Inches(0.1), Inches(0.32),
                 font_size=11, color=NAVY, align=PP_ALIGN.CENTER, bold=True)
        bx += bw + Inches(0.2)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — Demo Walkthrough
# ══════════════════════════════════════════════════════════════════════════════
def slide_demo(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    section_header(slide, "Demo Walkthrough",
                   "4 steps · ~15 seconds · coordinator stays in control")

    steps = [
        ("Step 1", "Submit\nCare Request",
         "Describe patient needs in\nJapanese or English.\n\n"
         "System identifies 介護保険\nservice codes + ward.\n\n"
         "Submits to 5-agent pipeline."),
        ("Step 2", "Watch\nAgents Work",
         "Pipeline animates in real time.\n\n"
         "5 agents run in sequence:\nintake → discovery → paperwork.\n\n"
         "Completes in ~15 seconds."),
        ("Step 3", "Human Reviews\n& Approves",
         "Coordinator sees ranked\nfacilities with match scores.\n\n"
         "All required forms pre-filled.\n\n"
         "Single click approves —\npipeline resumes."),
        ("Step 4", "Complete\nCare Plan",
         "4-week visit calendar\nauto-generated.\n\n"
         "Care plan scored 0–100\nwith risk alerts.\n\n"
         "LINE templates ready to send."),
    ]

    sw = Inches(2.88)
    sh = Inches(4.65)
    sg = Inches(0.27)
    sx = Inches(0.5)
    sy = Inches(1.35)

    for i, (step, title, body) in enumerate(steps):
        x = sx + i * (sw + sg)
        add_rect(slide, x, sy, sw, sh, WHITE,
                 line_color=MID_GRAY, line_width_pt=0.75)
        add_rect(slide, x, sy, sw, Inches(0.07), TEAL)

        # Step badge
        add_rect(slide, x + Inches(0.15), sy + Inches(0.14),
                 Inches(0.72), Inches(0.3), TEAL)
        add_text(slide, step,
                 x + Inches(0.15), sy + Inches(0.12),
                 Inches(0.72), Inches(0.3),
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_text(slide, title,
                 x + Inches(0.15), sy + Inches(0.55),
                 sw - Inches(0.3), Inches(0.75),
                 font_size=14, bold=True, color=NAVY)
        add_text(slide, body,
                 x + Inches(0.15), sy + Inches(1.35),
                 sw - Inches(0.3), Inches(3.1),
                 font_size=11, color=SLATE)

    # Time-saved banner
    add_rect(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.72),
             TEAL_LIGHT, line_color=TEAL, line_width_pt=1.0)
    add_rect(slide, Inches(0.5), Inches(6.2), Inches(0.07), Inches(0.72), TEAL)
    add_text(slide,
             "⏱  ~15 sec end-to-end  ·  3 forms pre-filled  ·  3 facilities ranked  ·  ~4 hours of paperwork saved",
             Inches(0.72), Inches(6.3), Inches(11.9), Inches(0.52),
             font_size=13, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — Use Case
# ══════════════════════════════════════════════════════════════════════════════
def slide_usecase(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    section_header(slide, "Live Use Case",
                   "Post-stroke patient · 世田谷区 · 介護度 2")

    # Patient profile card (left)
    add_rect(slide, Inches(0.5), Inches(1.35), Inches(5.6), Inches(2.65), WHITE,
             line_color=MID_GRAY, line_width_pt=0.75)
    add_rect(slide, Inches(0.5), Inches(1.35), Inches(5.6), Inches(0.07), TEAL)
    add_text(slide, "Patient Profile",
             Inches(0.65), Inches(1.43), Inches(5.3), Inches(0.4),
             font_size=13, bold=True, color=NAVY)

    info = [
        ("Name",      "田中 花子 (Tanaka Hanako)"),
        ("Age",       "78 years old"),
        ("Ward",      "世田谷区"),
        ("Care Level","介護度 2"),
        ("Condition", "Post-stroke (脳梗塞後遺症)"),
        ("Needs",     "週2回のデイサービスと週2回のリハビリが必要"),
    ]
    for j, (lbl, val) in enumerate(info):
        y = Inches(1.9) + j * Inches(0.33)
        add_text(slide, lbl + ":", Inches(0.65), y, Inches(1.3), Inches(0.3),
                 font_size=11, bold=True, color=TEAL)
        add_text(slide, val, Inches(2.05), y, Inches(3.9), Inches(0.3),
                 font_size=11, color=SLATE)

    # AI output card (right)
    add_rect(slide, Inches(6.5), Inches(1.35), Inches(6.3), Inches(2.65), WHITE,
             line_color=MID_GRAY, line_width_pt=0.75)
    add_rect(slide, Inches(6.5), Inches(1.35), Inches(6.3), Inches(0.07), TEAL)
    add_text(slide, "AI Pipeline Output  (15 sec)",
             Inches(6.65), Inches(1.43), Inches(6.0), Inches(0.4),
             font_size=13, bold=True, color=NAVY)

    outputs = [
        ("Service Codes",   "21 (通所介護)  ·  22 (通所リハビリ)"),
        ("Top Facility",    "世田谷デイサービスセンター  ·  Score: 88/100"),
        ("Forms Pre-filled","介護認定申請書, 通所介護利用申込書, 通所リハビリ申込書"),
        ("Schedule",        "Mon+Thu Day Service, Tue+Fri Rehab  ·  4 weeks"),
        ("Care Plan Score", "82 / 100  ·  Risk Level: Low ✅"),
    ]
    for j, (lbl, val) in enumerate(outputs):
        y = Inches(1.9) + j * Inches(0.34)
        add_text(slide, lbl + ":", Inches(6.65), y, Inches(2.1), Inches(0.3),
                 font_size=11, bold=True, color=TEAL)
        add_text(slide, val, Inches(8.85), y, Inches(3.8), Inches(0.3),
                 font_size=11, color=SLATE)

    # Without panel
    add_rect(slide, Inches(0.5), Inches(4.2), Inches(5.85), Inches(2.55),
             RGBColor(0xFF, 0xF5, 0xF5), line_color=MUTED_RED, line_width_pt=0.75)
    add_rect(slide, Inches(0.5), Inches(4.2), Inches(0.07), Inches(2.55), MUTED_RED)
    add_text(slide, "❌  Without Kaigo Navigator",
             Inches(0.72), Inches(4.28), Inches(5.5), Inches(0.42),
             font_size=13, bold=True, color=MUTED_RED)
    for j, line in enumerate([
        "4–6 weeks of manual form filling",
        "Family researches facilities independently",
        "No scoring — purely subjective decisions",
        "Reminders managed by phone calls",
    ]):
        add_text(slide, "·  " + line,
                 Inches(0.82), Inches(4.8) + j * Inches(0.38),
                 Inches(5.5), Inches(0.35), font_size=12, color=SLATE)

    # VS divider
    add_text(slide, "VS",
             Inches(6.15), Inches(5.2), Inches(0.7), Inches(0.55),
             font_size=22, bold=True, color=SUBTEXT, align=PP_ALIGN.CENTER)

    # With panel
    add_rect(slide, Inches(6.95), Inches(4.2), Inches(5.85), Inches(2.55),
             TEAL_LIGHT, line_color=TEAL, line_width_pt=0.75)
    add_rect(slide, Inches(6.95), Inches(4.2), Inches(0.07), Inches(2.55), TEAL)
    add_text(slide, "✅  With Kaigo Navigator",
             Inches(7.17), Inches(4.28), Inches(5.5), Inches(0.42),
             font_size=13, bold=True, color=TEAL_DARK)
    for j, line in enumerate([
        "15 seconds — fully automated pipeline",
        "3 facilities ranked with AI match scores",
        "3 forms pre-filled, ready to submit",
        "LINE reminders auto-generated",
    ]):
        add_text(slide, "·  " + line,
                 Inches(7.27), Inches(4.8) + j * Inches(0.38),
                 Inches(5.5), Inches(0.35), font_size=12, color=SLATE)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — What Makes This Different
# ══════════════════════════════════════════════════════════════════════════════
def slide_differentiators(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    section_header(slide, "What Makes This Different")

    points = [
        ("⏸",  "Human-in-the-Loop",
         "Pipeline halts before any submission. State persists across the HTTP pause. "
         "Coordinator reviews and approves — AI does the legwork, humans stay in control."),
        ("🌐", "Bilingual Throughout",
         "Every agent prompt, form field, ranked explanation, and dashboard label is in "
         "Japanese and English simultaneously. No switching, no translation step."),
        ("🛡", "Graceful Degradation",
         "Every LLM call has a deterministic rule-based fallback. If Groq is down, "
         "the pipeline completes with heuristic output. It never crashes."),
        ("📊", "Live Government Data",
         "Scraper pulls directly from MHLW's 介護サービス情報公表システム — Japan's official "
         "care facility registry — with ward-level filtering across all 23 Tokyo wards."),
    ]

    for i, (icon, title, body) in enumerate(points):
        y = Inches(1.3) + i * Inches(1.35)
        add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(1.18), WHITE,
                 line_color=MID_GRAY, line_width_pt=0.75)
        add_rect(slide, Inches(0.5), y, Inches(0.07), Inches(1.18), TEAL)

        # Icon circle (simulated with small rect)
        add_rect(slide, Inches(0.72), y + Inches(0.3),
                 Inches(0.5), Inches(0.5), TEAL_LIGHT)
        add_text(slide, icon,
                 Inches(0.72), y + Inches(0.28), Inches(0.5), Inches(0.5),
                 font_size=16, align=PP_ALIGN.CENTER)

        add_text(slide, title,
                 Inches(1.42), y + Inches(0.1),
                 Inches(10.8), Inches(0.42),
                 font_size=15, bold=True, color=NAVY)
        add_text(slide, body,
                 Inches(1.42), y + Inches(0.55),
                 Inches(10.8), Inches(0.52),
                 font_size=12, color=SLATE)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — Tech Stack
# ══════════════════════════════════════════════════════════════════════════════
def slide_tech(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    section_header(slide, "Tech Stack")

    rows = [
        ("UI",            "Next.js 16 + Tailwind CSS",     "3-page care coordinator interface"),
        ("API",           "FastAPI + Pydantic v2",          "REST layer, async, auto-documented"),
        ("Orchestration", "LangGraph 0.6",                 "Stateful multi-agent graph + HITL interrupt"),
        ("LLM",           "Groq — Llama 3.3 70B / 3.1 8B","Sub-second inference, bilingual JP+EN"),
        ("Vector DB",     "Pinecone",                       "Ward + service code metadata filtering"),
        ("Embeddings",    "multilingual-e5-large",          "JP + EN in the same vector space, local"),
        ("Observability", "Langfuse",                       "Full trace on every agent step"),
        ("State",         "SQLite + langgraph-checkpoint",  "HITL pause/resume across HTTP requests"),
    ]

    col_w = [Inches(2.1), Inches(3.7), Inches(6.0)]
    headers = ["Layer", "Technology", "Role"]
    hy = Inches(1.3)

    for ci, (hdr, cw) in enumerate(zip(headers, col_w)):
        cx = Inches(0.5) + sum(col_w[:ci])
        add_rect(slide, cx, hy, cw, Inches(0.4), NAVY)
        add_text(slide, hdr,
                 cx + Inches(0.12), hy + Inches(0.05),
                 cw - Inches(0.24), Inches(0.32),
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for ri, (layer, tech, role) in enumerate(rows):
        ry = Inches(1.72) + ri * Inches(0.57)
        bg = WHITE if ri % 2 == 0 else SOFT_GRAY
        for ci, (val, cw) in enumerate(zip([layer, tech, role], col_w)):
            cx = Inches(0.5) + sum(col_w[:ci])
            add_rect(slide, cx, ry, cw, Inches(0.55), bg,
                     line_color=MID_GRAY, line_width_pt=0.3)
            clr = TEAL_DARK if ci == 0 else (NAVY if ci == 1 else SLATE)
            add_text(slide, val,
                     cx + Inches(0.12), ry + Inches(0.1),
                     cw - Inches(0.24), Inches(0.38),
                     font_size=12, bold=(ci < 2), color=clr)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — Project Status
# ══════════════════════════════════════════════════════════════════════════════
def slide_status(prs):
    slide = blank_slide(prs)
    fill_bg(slide)
    section_header(slide, "Project Status & Roadmap")

    phases = [
        ("✅", "Phase 1 — Complete",
         "RAG pipeline  ·  Service Discovery  ·  Langfuse observability",
         TEAL, TEAL_LIGHT),
        ("✅", "Phase 2 — Complete",
         "Paperwork Agent  ·  Human-in-the-loop approval gate",
         TEAL, TEAL_LIGHT),
        ("✅", "Phase 3 — Complete",
         "Scheduling Agent  ·  Monitoring Agent  ·  Next.js UI",
         TEAL, TEAL_LIGHT),
        ("🔄", "Phase 4 — In Progress",
         "Eval harness (50 synthetic cases)  ·  Cloud deployment",
         AMBER, RGBColor(0xFF, 0xF8, 0xE6)),
    ]

    for i, (icon, phase, scope, bar_color, bg_color) in enumerate(phases):
        y = Inches(1.3) + i * Inches(1.3)
        add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(1.1), bg_color,
                 line_color=bar_color, line_width_pt=0.75)
        add_rect(slide, Inches(0.5), y, Inches(0.07), Inches(1.1), bar_color)

        add_text(slide, icon + "  " + phase,
                 Inches(0.75), y + Inches(0.1),
                 Inches(6), Inches(0.45),
                 font_size=16, bold=True,
                 color=NAVY if bar_color == TEAL else RGBColor(0x92, 0x40, 0x00))
        add_text(slide, scope,
                 Inches(0.75), y + Inches(0.58),
                 Inches(11.4), Inches(0.4),
                 font_size=13, color=SLATE)

    add_rect(slide, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.02), MID_GRAY)
    add_text(slide, "Phases 1–3 fully functional and demo-ready today.",
             Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.38),
             font_size=13, color=SUBTEXT, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — Call to Action
# ══════════════════════════════════════════════════════════════════════════════
def slide_cta(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)

    # Hero navy block
    add_rect(slide, 0, 0, SLIDE_W, Inches(3.2), NAVY)
    add_rect(slide, 0, Inches(3.2), SLIDE_W, Inches(0.07), TEAL)

    add_text(slide, "Let's Build This Together",
             Inches(1), Inches(0.55), Inches(11.3), Inches(1.0),
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide,
             "Kaigo Navigator is built and running.\n"
             "Phases 1–3 are complete and demo-ready today.",
             Inches(1.5), Inches(1.65), Inches(10.3), Inches(1.1),
             font_size=17, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Three CTA cards
    ctas = [
        ("🖥", "Live Demo",
         "Submit a real care request\nand watch 5 agents run live"),
        ("📁", "GitHub Repo",
         "Full source code, README,\nand interactive API docs"),
        ("📊", "Langfuse Traces",
         "Every agent step traced\nand observable in real time"),
    ]
    cw = Inches(3.6)
    ch = Inches(2.4)
    cg = Inches(0.55)
    cx = Inches(0.92)
    cy = Inches(3.65)

    for icon, title, desc in ctas:
        add_rect(slide, cx, cy, cw, ch, WHITE,
                 line_color=MID_GRAY, line_width_pt=0.75)
        add_rect(slide, cx, cy, cw, Inches(0.07), TEAL)
        add_text(slide, icon,
                 cx, cy + Inches(0.2), cw, Inches(0.55),
                 font_size=26, align=PP_ALIGN.CENTER)
        add_text(slide, title,
                 cx, cy + Inches(0.85), cw, Inches(0.45),
                 font_size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, desc,
                 cx + Inches(0.15), cy + Inches(1.35),
                 cw - Inches(0.3), Inches(0.85),
                 font_size=12, color=SLATE, align=PP_ALIGN.CENTER)
        cx += cw + cg

    # Footer
    add_rect(slide, 0, SLIDE_H - Inches(0.48), SLIDE_W, Inches(0.48), SOFT_GRAY)
    add_text(slide, "Built by Ganesh Suni Jathar",
             Inches(0.5), SLIDE_H - Inches(0.44), Inches(6), Inches(0.38),
             font_size=11, color=SUBTEXT)
    add_text(slide, "github.com/your-handle/kaigo-navigator",
             Inches(7.0), SLIDE_H - Inches(0.44), Inches(6.0), Inches(0.38),
             font_size=11, color=TEAL, italic=True, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════
def main():
    prs = new_prs()

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_demo(prs)
    slide_usecase(prs)
    slide_differentiators(prs)
    slide_tech(prs)
    slide_status(prs)
    slide_cta(prs)

    os.makedirs("docs", exist_ok=True)
    out = "docs/kaigo_navigator_demo.pptx"
    prs.save(out)
    print(f"✅  Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
